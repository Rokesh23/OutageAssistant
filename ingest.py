import os
import pickle
import re
import faiss
import numpy as np
import ollama

from config import *

import docx
import openpyxl
from pptx import Presentation

# ==========================================================
# ADVANCED & COMPREHENSIVE DOCUMENT READERS
# ==========================================================

def read_docx(file_path):
    """
    Reads ALL text in exact structural order by extracting text from 
    paragraphs, tables, merged cells, and XML text nodes without skipping.
    """
    doc = docx.Document(file_path)
    full_text = []

    # Iterate through all direct child elements of the document body (in exact reading order)
    for element in doc.element.body:
        # Paragraph element
        if element.tag.endswith('p'):
            p = docx.text.paragraph.Paragraph(element, doc)
            if p.text.strip():
                full_text.append(p.text.strip())

        # Table element
        elif element.tag.endswith('tbl'):
            table = docx.table.Table(element, doc)
            for row in table.rows:
                row_cells = []
                for cell in row.cells:
                    # Collect text using itertext to catch all nested XML text nodes
                    cell_str = " ".join("".join(cell._tc.itertext()).split())
                    if cell_str and cell_str not in row_cells:
                        row_cells.append(cell_str)
                if row_cells:
                    full_text.append(" | ".join(row_cells))

    # Fallback to full body text if structural parsing yielded no text
    if not full_text:
        text_content = "".join(doc.element.body.itertext()).strip()
        if text_content:
            full_text.append(text_content)

    return "\n\n".join(full_text)


def read_xlsx(file_path):
    wb = openpyxl.load_workbook(file_path, data_only=True)
    full_text = []
    for sheet in wb.sheetnames:
        ws = wb[sheet]
        for row in ws.iter_rows(values_only=True):
            row_vals = [str(val).strip() for val in row if val is not None and str(val).strip() != ""]
            if row_vals:
                full_text.append(" | ".join(row_vals))
    return "\n\n".join(full_text)


def read_pptx(file_path):
    prs = Presentation(file_path)
    full_text = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                full_text.append(shape.text.strip())
    return "\n\n".join(full_text)


def read_file(file_path):
    ext = os.path.splitext(file_path)[1].lower()
    try:
        if ext == ".docx":
            return read_docx(file_path)
        elif ext == ".xlsx":
            return read_xlsx(file_path)
        elif ext == ".pptx":
            return read_pptx(file_path)
        elif ext in [".txt", ".md"]:
            with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
                return f.read()
    except Exception as e:
        print(f"Error reading {file_path}: {e}")
        return ""
    return ""


# ==========================================================
# INCIDENT ID PARSER
# ==========================================================

def parse_incident_id(filename, text):
    match = re.search(r"(INC\d+|QA-INC-\d+)", filename.upper())
    if match:
        return match.group(0)
    match_text = re.search(r"(INC\d+|QA-INC-\d+)", text.upper())
    if match_text:
        return match_text.group(0)
    return "GENERAL"


# ==========================================================
# SMART PARAGRAPH-AWARE CHUNKER
# ==========================================================

def chunk_text(text, max_chars=1000):
    paragraphs = text.split("\n\n")
    chunks = []
    current_chunk = ""

    for p in paragraphs:
        p_clean = p.strip()
        if not p_clean:
            continue
        if len(current_chunk) + len(p_clean) < max_chars:
            current_chunk += p_clean + "\n\n"
        else:
            if current_chunk:
                chunks.append(current_chunk.strip())
            current_chunk = p_clean + "\n\n"

    if current_chunk.strip():
        chunks.append(current_chunk.strip())

    return chunks if chunks else [text[:max_chars]]


# ==========================================================
# MAIN INGESTION PIPELINE
# ==========================================================

def run_ingestion():
    documents_dir = DOCUMENTS_FOLDER if 'DOCUMENTS_FOLDER' in globals() else "./documents"
    metadata = []
    embeddings_list = []

    if not os.path.exists(documents_dir):
        print(f"Error: Directory '{documents_dir}' does not exist.")
        return

    valid_extensions = (".docx", ".xlsx", ".pptx", ".txt", ".md")
    files = [f for f in os.listdir(documents_dir) if f.lower().endswith(valid_extensions)]
    
    print(f"Starting ingestion process for {len(files)} document(s)...\n")

    for file_name in files:
        file_path = os.path.join(documents_dir, file_name)
        text = read_file(file_path)

        if not text.strip():
            print(f"⚠️ Skipping empty or unreadable file: {file_name}")
            continue

        incident_id = parse_incident_id(file_name, text)
        chunks = chunk_text(text)

        print(f"📄 Indexed: {file_name} -> Incident ID: {incident_id} ({len(chunks)} chunks)")

        for chunk in chunks:
            response = ollama.embed(model=EMBEDDING_MODEL, input=chunk)
            emb = response["embeddings"][0]
            
            embeddings_list.append(emb)
            metadata.append({
                "source": file_name,
                "incident": incident_id,
                "text": chunk
            })

    if not embeddings_list:
        print("❌ No text content extracted from documents.")
        return

    embeddings_np = np.array(embeddings_list, dtype="float32")
    dimension = embeddings_np.shape[1]

    index = faiss.IndexFlatL2(dimension)
    index.add(embeddings_np)

    if not os.path.exists(INDEX_FOLDER):
        os.makedirs(INDEX_FOLDER)

    faiss.write_index(index, f"{INDEX_FOLDER}/index.faiss")
    with open(f"{INDEX_FOLDER}/metadata.pkl", "wb") as f:
        pickle.dump(metadata, f)

    print(f"\n✅ Ingestion complete! Successfully indexed {len(files)} files into {len(metadata)} chunks.")


# Explicit execution entrypoint
if __name__ == "__main__":
    run_ingestion()